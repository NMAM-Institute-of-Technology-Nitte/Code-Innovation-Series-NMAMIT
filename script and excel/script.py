# import Presentation class 
# from rootx library 
from pptx import Presentation      
import xlrd
from pptx.util import Inches, Pt 
 
loc = ("./Hackathon Admin Panel  INCUBATEIND.xlsx")
 
wb = xlrd.open_workbook(loc)
sheet = wb.sheet_by_index(0)
 
for i in range(2,81):
    root = Presentation() 
    
    first_slide_layout = root.slide_layouts[0]  
    slide = root.slides.add_slide(first_slide_layout) 
    slide.shapes.title.text = sheet.cell_value(i, 2)
    s="Team Lead: "+sheet.cell_value(i, 4)+"\nTeam Lead Email: "+sheet.cell_value(i, 3)+"\nMembers:\n"+sheet.cell_value(i, 11)+"\n"+sheet.cell_value(i, 14)+"\n"+sheet.cell_value(i, 17)
    slide.placeholders[1].text = s

    
    #blank_slide_layout = root.slide_layouts[6]  
    #slide = root.slides.add_slide(blank_slide_layout) 
    #left = top = width = height = Inches(1)  
    #txBox = slide.shapes.add_textbox(left, top, 0.5, 100) 
    #tf = txBox.text_frame 
    #p = tf.add_paragraph()  
    #p.text = sheet.cell_value(i, 20) 
    #p.font.bold = True
    #p.font.size = Pt(12) 


    second_slide_layout = root.slide_layouts[0]  
    slide = root.slides.add_slide(second_slide_layout) 
    text_ph1= slide.placeholders[0]
    text_ph1.text = sheet.cell_value(i, 20)
    try:
        font = text_ph1.text_frame.paragraphs[0].runs[0].font
        font.size = Pt(16)
        font.bold = True
    except:
        pass

    second_slide_layout = root.slide_layouts[0]  
    slide = root.slides.add_slide(second_slide_layout) 
    text_ph2= slide.placeholders[0]
    text_ph2.text = sheet.cell_value(i, 21)
    
    try:
        font = text_ph2.text_frame.paragraphs[0].runs[0].font
        font.size = Pt(16)
        font.bold = True
    except:
        pass

    root.save(sheet.cell_value(i, 2)+".pptx") 
  

print("done") 